from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


OUT = Path(r"D:\plantAssistant\Пояснительная записка\Защитная презентация.pptx")

GREEN = RGBColor(22, 132, 58)
GREEN_DARK = RGBColor(12, 94, 42)
GREEN_SOFT = RGBColor(223, 243, 229)
BG = RGBColor(245, 248, 243)
TEXT = RGBColor(25, 35, 25)
MUTED = RGBColor(90, 102, 92)
BORDER = RGBColor(204, 220, 205)
WHITE = RGBColor(255, 255, 255)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_top_band(slide, title, subtitle=None):
    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.0)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = GREEN
    band.line.color.rgb = GREEN

    tx = slide.shapes.add_textbox(Inches(0.55), Inches(0.22), Inches(9.8), Inches(0.45))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos Display"
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = WHITE

    if subtitle:
        st = slide.shapes.add_textbox(Inches(0.58), Inches(0.64), Inches(8.5), Inches(0.22))
        tf = st.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = subtitle
        r.font.name = "Aptos"
        r.font.size = Pt(10.5)
        r.font.color.rgb = WHITE


def add_text_box(slide, x, y, w, h, lines, font_size=20, color=TEXT, bold_first=False):
    shape = slide.shapes.add_textbox(x, y, w, h)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(4)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(6)
        r = p.add_run()
        r.text = line
        r.font.name = "Aptos"
        r.font.size = Pt(font_size)
        r.font.color.rgb = color
        if bold_first and idx == 0:
            r.font.bold = True
    return shape


def add_bullets(slide, x, y, w, h, bullets, font_size=19):
    shape = slide.shapes.add_textbox(x, y, w, h)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(6)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for idx, line in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(8)
        p.bullet = True
        r = p.add_run()
        r.text = line
        r.font.name = "Aptos"
        r.font.size = Pt(font_size)
        r.font.color.rgb = TEXT
    return shape


def add_placeholder(slide, x, y, w, h, title, body):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = BORDER
    card.line.width = Pt(1.25)

    hdr = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, Inches(0.42))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = GREEN_SOFT
    hdr.line.color.rgb = GREEN_SOFT

    add_text_box(slide, x + Inches(0.12), y + Inches(0.05), w - Inches(0.24), Inches(0.25), [title], 12, GREEN_DARK, True)
    add_text_box(slide, x + Inches(0.16), y + Inches(0.58), w - Inches(0.32), h - Inches(0.72), [body], 15, MUTED)


def add_table(slide, x, y, w, h, headers, rows):
    table = slide.shapes.add_table(len(rows) + 1, len(headers), x, y, w, h).table
    table.first_row = True
    col_w = int(w / len(headers))
    for i in range(len(headers)):
        table.columns[i].width = col_w
    for i, head in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = GREEN
        cell.text = head
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            r.font.name = "Aptos"
            r.font.size = Pt(12)
            r.font.bold = True
            r.font.color.rgb = WHITE
    for r_idx, row in enumerate(rows, start=1):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r_idx % 2 else GREEN_SOFT
            cell.text = value
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            for run in p.runs:
                run.font.name = "Aptos"
                run.font.size = Pt(11)
                run.font.color.rgb = TEXT


def title_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.75))
    band.fill.solid()
    band.fill.fore_color.rgb = GREEN
    band.line.color.rgb = GREEN

    add_text_box(
        slide,
        Inches(0.75),
        Inches(1.1),
        Inches(11.5),
        Inches(1.8),
        ["Разработка веб-приложения клуба любителей растений Plant Assistant"],
        27,
        TEXT,
        True,
    )
    add_text_box(
        slide,
        Inches(0.78),
        Inches(3.15),
        Inches(8.4),
        Inches(1.0),
        [
            "Выпускная квалификационная работа",
            "Направление: разработка информационной системы для учета растений и планирования ухода",
        ],
        18,
        MUTED,
    )
    add_placeholder(
        slide,
        Inches(8.95),
        Inches(2.5),
        Inches(3.45),
        Inches(2.1),
        "Титульный блок",
        "Вставить ФИО, группу, руководителя и год защиты."
    )


def content_slide(title, bullets=None, left_box=None, right_placeholder=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide, title)
    if bullets:
        add_bullets(slide, Inches(0.55), Inches(1.25), Inches(5.75), Inches(5.5), bullets)
    if left_box:
        add_text_box(slide, Inches(0.55), Inches(1.35), Inches(5.7), Inches(5.3), left_box, 18)
    if right_placeholder:
        add_placeholder(slide, Inches(6.7), Inches(1.35), Inches(5.95), Inches(4.9), right_placeholder[0], right_placeholder[1])
    return slide


title_slide()

content_slide(
    "Актуальность, цель, задачи",
    bullets=[
        "Уход за комнатными растениями требует регулярности и систематизации.",
        "Данные о растениях часто хранятся разрозненно.",
        "Цель: разработка веб-приложения для хранения, обработки и визуализации данных о растениях и уходе.",
        "Задачи: анализ, проектирование, реализация backend и frontend, тестирование.",
    ],
    right_placeholder=("Схема проблемы", "Вставить схему: заметки, календари и мессенджеры -> отсутствие единой системы -> необходимость Plant Assistant."),
)

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_top_band(slide, "Основные технологии backend")
add_table(
    slide,
    Inches(0.6),
    Inches(1.35),
    Inches(6.2),
    Inches(4.5),
    ["Технология", "Назначение"],
    [
        ["PHP 8.3", "серверная логика"],
        ["Laravel 13", "REST API и бизнес-логика"],
        ["PostgreSQL", "хранение данных"],
        ["Eloquent ORM", "модели и связи"],
        ["Sanctum", "авторизация"],
        ["PHPUnit", "тестирование"],
    ],
)
add_placeholder(slide, Inches(7.2), Inches(1.35), Inches(5.45), Inches(4.9), "Комментарий", "На защите кратко пояснить, почему именно этот стек подходит для API-приложения.")

content_slide(
    "PHP и Laravel",
    bullets=[
        "PHP используется как основной язык backend.",
        "Laravel является каркасом REST API.",
        "Используются контроллеры, модели, middleware, policies.",
        "Form Requests отвечают за валидацию, API Resources - за формат ответов.",
    ],
    right_placeholder=("Схема Laravel", "Вставить схему: HTTP-запрос -> route -> controller -> service/model -> JSON-ответ."),
)

content_slide(
    "База данных, ORM, миграции и сидеры",
    bullets=[
        "PostgreSQL выбрана как реляционная СУБД.",
        "Eloquent ORM упрощает работу с сущностями и связями.",
        "Миграции позволяют управлять структурой базы данных.",
        "Сидеры используются для начального заполнения ролями, пользователями и тестовыми данными.",
    ],
    right_placeholder=("ER-схема", "Вставить ER-схему: users, roles, plants, rooms, care_settings, care_logs, tips, likes, follows, reports."),
)

content_slide(
    "Laravel Sanctum и система авторизации",
    bullets=[
        "Для авторизации используется Laravel Sanctum.",
        "После входа пользователь получает Bearer token.",
        "Доступ разграничивается по ролям user и admin.",
        "Права дополнительно контролируются через middleware и policies.",
        "Корректный термин для защиты: Bearer token, а не JWT.",
    ],
    right_placeholder=("Схема авторизации", "Вставить схему: login -> токен -> Authorization: Bearer <token> -> проверка доступа."),
)

content_slide(
    "Docker и RoadRunner",
    bullets=[
        "Docker обеспечивает воспроизводимое развертывание проекта.",
        "Контейнеризация уменьшает зависимость от локального окружения.",
        "RoadRunner использует долгоживущие worker-процессы.",
        "Laravel загружается один раз и обслуживает множество запросов.",
        "Это уменьшает накладные расходы на повторную инициализацию приложения.",
    ],
    right_placeholder=("Сравнение запуска", "Вставить схему сравнения: обычный PHP-запуск на каждый запрос и модель RoadRunner с worker-процессами."),
)

content_slide(
    "Vue.js и Vite",
    bullets=[
        "Frontend построен на Vue.js как SPA-приложение.",
        "Компонентный подход упрощает поддержку и переиспользование UI.",
        "Vite используется для dev-сервера и production-сборки.",
        "Маршрутизация реализована через Vue Router.",
    ],
    right_placeholder=("Скриншот интерфейса", "Вставить скриншот ленты, задач или формы растения."),
)

content_slide(
    "Pinia",
    bullets=[
        "Pinia используется для управления состоянием приложения.",
        "Отдельные store созданы для auth, plants, tasks, social и admin.",
        "Состояние отделено от визуальных компонентов.",
        "Это уменьшает дублирование логики и упрощает сопровождение.",
    ],
    right_placeholder=("Схема состояния", "Вставить схему: UI-компоненты -> Pinia store -> API client -> backend."),
)

content_slide(
    "Единый API-клиент",
    bullets=[
        "Во frontend используется единый API-клиент.",
        "Он централизует HTTP-запросы к backend.",
        "Поддерживает JSON, FormData и Bearer token.",
        "Единообразно обрабатывает ошибки и ответы API.",
    ],
    right_placeholder=("Схема взаимодействия", "Вставить схему: page/component -> apiClient -> REST API."),
)

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_top_band(slide, "VeeValidate, Zod и другие пакеты frontend")
add_table(
    slide,
    Inches(0.6),
    Inches(1.35),
    Inches(6.55),
    Inches(4.9),
    ["Пакет", "Применение"],
    [
        ["VeeValidate", "формы ввода"],
        ["Zod", "схемы проверки данных"],
        ["Chart.js + vue-chartjs", "графики и статистика"],
        ["lucide-vue-next", "иконки интерфейса"],
        ["vue-sonner", "уведомления"],
        ["@vueuse/core", "адаптивность и composable"],
    ],
)
add_placeholder(slide, Inches(7.55), Inches(1.35), Inches(5.1), Inches(4.9), "Где используются", "На защите привести примеры: форма растения, профиль, графики статистики, иконки интерфейса и уведомления.")

content_slide(
    "Тестирование frontend и backend",
    bullets=[
        "Backend тестируется через PHPUnit и OpenAPI-контракт.",
        "Проверяются API-маршруты, авторизация и административные сценарии.",
        "Frontend проходит сборку, smoke-check и статическую проверку.",
        "Для контроля качества используются ESLint, Prettier и Laravel Pint.",
    ],
    right_placeholder=("Схема тестирования", "Вставить схему: backend tests + frontend checks + ручная демонстрация."),
)

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_top_band(slide, "Функционал и демонстрация")
add_bullets(
    slide,
    Inches(0.6),
    Inches(1.3),
    Inches(4.7),
    Inches(4.8),
    [
        "Регистрация и вход",
        "Добавление растения",
        "Настройка ухода",
        "Список задач",
        "Публичная лента",
        "Административная панель",
    ],
    19,
)
box_w = Inches(3.15)
box_h = Inches(1.65)
start_x = Inches(5.7)
start_y = Inches(1.45)
gap_x = Inches(0.25)
gap_y = Inches(0.25)
labels = ["Лента", "Форма растения", "Задачи", "Админ-панель"]
for idx, label in enumerate(labels):
    row = idx // 2
    col = idx % 2
    x = start_x + col * (box_w + gap_x)
    y = start_y + row * (box_h + gap_y)
    add_placeholder(slide, x, y, box_w, box_h, label, f"Вставить скриншот: {label.lower()}.")

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUT))
print(OUT)
